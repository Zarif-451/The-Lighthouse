import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Modern Light Academic Palette
    BG_COLOR = RGBColor(248, 250, 252)       # Light Slate #F8FAFC
    TEXT_MAIN = RGBColor(15, 23, 42)         # Dark Slate #0F172A
    TEXT_MUTED = RGBColor(71, 85, 105)       # Muted Slate #475569
    ACCENT_BLUE = RGBColor(37, 99, 235)      # Primary Blue #2563EB
    ACCENT_TEAL = RGBColor(13, 148, 136)     # Primary Teal #0D9488
    BLUE_BG = RGBColor(239, 246, 255)        # Light Blue #EFF6FF
    BORDER_BLUE = RGBColor(191, 219, 254)    # Border Blue #BFDBFE
    
    blank_layout = prs.slide_layouts[6]
    
    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_BLUE
        bar.line.fill.background()
        return slide

    def add_header(slide, tag_text, title_text):
        add_bg(slide)
        
        tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.3))
        p_tag = tb_tag.text_frame.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.name = 'Segoe UI'
        p_tag.font.size = Pt(10.5)
        p_tag.font.bold = True
        p_tag.font.color.rgb = ACCENT_BLUE
        
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.55))
        p_title = tb_title.text_frame.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Segoe UI'
        p_title.font.size = Pt(25)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_MAIN

    def add_single_image_slide(tag, title, points, img_path, caption, align="left_text"):
        s = prs.slides.add_slide(blank_layout)
        add_header(s, tag, title)
        
        max_w = Inches(7.6)
        max_h = Inches(4.7)
        img_top = Inches(1.5)
        text_w = Inches(4.3)
        text_top = Inches(1.6)
        
        if align == "left_text":
            text_left = Inches(0.8)
            img_left = Inches(5.1)
        else:
            text_left = Inches(8.2)
            img_left = Inches(0.8)
            
        tb = s.shapes.add_textbox(text_left, text_top, text_w, Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        for pt in points:
            p = tf.add_paragraph() if len(tf.paragraphs) > 0 and tf.paragraphs[0].text else tf.paragraphs[0]
            p.text = "●  " + pt
            p.font.name = 'Segoe UI'
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MAIN
            p.space_after = Pt(16)
            
        if os.path.exists(img_path):
            try:
                pic = s.shapes.add_picture(img_path, img_left, img_top, width=max_w)
                if pic.height > max_h:
                    pic.height = max_h
                    pic.left = img_left + int((max_w - pic.width) / 2)
                    
                pic.line.color.rgb = BORDER_BLUE
                pic.line.width = Pt(1)
                
                tb_cap = s.shapes.add_textbox(img_left, pic.top + pic.height + Inches(0.04), pic.width, Inches(0.3))
                p_cap = tb_cap.text_frame.paragraphs[0]
                p_cap.text = caption
                p_cap.alignment = PP_ALIGN.CENTER
                p_cap.font.name = 'Segoe UI'
                p_cap.font.size = Pt(10)
                p_cap.font.italic = True
                p_cap.font.color.rgb = TEXT_MUTED
            except Exception as e:
                print(f"Error loading image {img_path}: {e}")

    def add_dual_image_slide(tag, title, card1_info, card2_info):
        s = prs.slides.add_slide(blank_layout)
        add_header(s, tag, title)
        
        card_w = Inches(5.6)
        card_h = Inches(5.2)
        top_y = Inches(1.5)
        
        for idx, (img_p, lbl, desc) in enumerate([card1_info, card2_info]):
            left_x = Inches(0.8) if idx == 0 else Inches(6.9)
            
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top_y, card_w, card_h)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            box.line.color.rgb = BORDER_BLUE
            
            if os.path.exists(img_p):
                try:
                    pic = s.shapes.add_picture(img_p, left_x + Inches(0.2), top_y + Inches(0.2), width=Inches(5.2))
                    if pic.height > Inches(3.2):
                        pic.height = Inches(3.2)
                    pic.line.color.rgb = BORDER_BLUE
                    pic.line.width = Pt(1)
                except Exception as e:
                    print(f"Error loading picture: {e}")
            
            tb = s.shapes.add_textbox(left_x + Inches(0.2), top_y + Inches(3.5), Inches(5.2), Inches(1.5))
            tf = tb.text_frame
            tf.word_wrap = True
            
            p1 = tf.paragraphs[0]
            p1.text = lbl
            p1.font.name = 'Segoe UI'
            p1.font.size = Pt(14)
            p1.font.bold = True
            p1.font.color.rgb = ACCENT_BLUE
            p1.space_after = Pt(4)
            
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.name = 'Segoe UI'
            p2.font.size = Pt(11)
            p2.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 1: COVER SLIDE & TEAM INTRO
    # =============================================================
    s1 = prs.slides.add_slide(blank_layout)
    add_bg(s1)
    
    logo = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.16), Inches(0.4), Inches(1.0), Inches(1.0))
    logo.fill.solid()
    logo.fill.fore_color.rgb = ACCENT_BLUE
    logo.line.fill.background()
    
    tb_tag = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(0.4))
    p = tb_tag.text_frame.paragraphs[0]
    p.text = "SOFTWARE DEVELOPMENT WITH JAVA (SESSIONAL) — CSE-202"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    tb_title = s1.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(11.333), Inches(0.9))
    p = tb_title.text_frame.paragraphs[0]
    p.text = "Project Lighthouse"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    tb_sub = s1.shapes.add_textbox(Inches(1.0), Inches(2.95), Inches(11.333), Inches(0.5))
    p = tb_sub.text_frame.paragraphs[0]
    p.text = "Wellbeing and Self-Reflection Web Platform"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = ACCENT_TEAL
    
    intro_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.6), Inches(10.933), Inches(3.2))
    intro_card.fill.solid()
    intro_card.fill.fore_color.rgb = BLUE_BG
    intro_card.line.color.rgb = BORDER_BLUE
    
    tb_intro = s1.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.333), Inches(2.8))
    tf_i = tb_intro.text_frame
    tf_i.word_wrap = True
    
    info_lines = [
        ("Course:", "Software Development with JAVA (Sessional) (CSE-202)"),
        ("Department:", "Department of Computer Science & Engineering, CUET"),
        ("Development Team:", "Muhammad Zarif Rahman (2304039)  |  Md. Al-Saiban (2304041)  |  Md. Irfan (2304066)")
    ]
    for label, val in info_lines:
        p = tf_i.add_paragraph() if len(tf_i.paragraphs) > 0 and tf_i.paragraphs[0].text else tf_i.paragraphs[0]
        p.text = f"{label} "
        p.font.name = 'Segoe UI'
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        run = p.add_run()
        run.text = val
        run.font.bold = False
        run.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(12)

    # SLIDE 2: THE PROBLEM
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "THE PROBLEM", "Why Does Mental Wellbeing Go Untracked?")
    probs = [
        ("No Structured Habit", "People want to reflect daily but have no guided framework to do it consistently."),
        ("Fear of Privacy Invasion", "Most organizational wellness tools are rejected out of fear reflections are read by managers."),
        ("No Longitudinal Insight", "Standalone mood apps capture a single moment but never correlate sleep, stress, and activity over time."),
        ("No Gentle Intervention", "Admins or counsellors have no visibility into who may need a gentle check-in until it's too late.")
    ]
    for idx, (h, d) in enumerate(probs):
        col = idx % 2
        row = idx // 2
        lx = Inches(0.8) if col == 0 else Inches(6.9)
        ty = Inches(1.6) if row == 0 else Inches(4.3)
        
        c = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, ty, Inches(5.6), Inches(2.3))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(255, 255, 255)
        c.line.color.rgb = BORDER_BLUE
        
        tb = s2.shapes.add_textbox(lx + Inches(0.3), ty + Inches(0.2), Inches(5.0), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = h
        p1.font.name = 'Segoe UI'
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_BLUE
        p1.space_after = Pt(6)
        p2 = tf.add_paragraph()
        p2.text = d
        p2.font.name = 'Segoe UI'
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED

    # SLIDE 3: APP INTRODUCTION
    add_single_image_slide("INTRODUCTION", "What is Lighthouse?", [
        "Core Mission: A guided daily wellbeing platform pairing a 13-step reflection journey with strict user privacy.",
        "Privacy Moat: Built on Supabase Row Level Security (RLS). Journal entries are cryptographically scoped to the user.",
        "Longitudinal Analytics: Interactive Chart.js visualisations reveal hidden sleep-mood-energy correlations over time.",
        "Proactive Admin Care: Heuristic flagging + soft nudges allow admins to reach out compassionately without invading privacy."
    ], "section3_2_dashboard.png", "Lighthouse Wellbeing Platform Overview", "left_text")

    # SLIDE 4: LANDING & AUTH
    add_dual_image_slide("AUTHENTICATION", "Landing Page & User Sign-Up", 
        ("section3_1_landing.png", "Marketing + Auth Page", "Single page hosting Login, Sign Up, and Admin tabs. New visitors see product pitch; returning users access auth."),
        ("section3_1_signup.png", "Extended Sign-Up Profile", "Collects extended profile fields - display name, avatar, occupation, interests, DOB - powering personalized AI context.")
    )

    # SLIDE 5: THE DAILY JOURNEY
    add_single_image_slide("CORE LOOP", "The Daily Journey — 13 Steps", [
        "Sequential Flow: Users progress through a 13-step structured sequence each day.",
        "Varied Micro-Activities: Check-ins, coping scenario assessments, cognitive visual challenges, and freeform journal writing.",
        "Progress Tracking: Clear step indicator ('Step X of 13') keeps daily reflection lightweight and engaging."
    ], "section3_2_dashboard.png", "Daily Journey & Progress Tracker", "left_text")

    # SLIDE 6: CHECK-IN & SCENARIO
    add_dual_image_slide("GUIDED JOURNEY — STEPS 1 & 3, 5, 7, 9, 11", "Daily Check-in & Scenario Assessment",
        ("section3_3_checkin.png", "Step 1 — Daily Check-in", "Users log sleep hours, mood (1-5), energy, productivity, and physical activity via smooth sliders in under 60 seconds."),
        ("section3_5_scenario.png", "Steps 3, 5, 7, 9, 11 — Scenario Assessment (x5)", "Five scenario rounds present situational stories from an 80+ item bank to probe coping patterns without clinical labels.")
    )

    # SLIDE 7: REFLECTION JOURNAL
    add_single_image_slide("GUIDED JOURNEY — STEPS 6 & 13", "Visual Reflection & Private Journal", [
        "Step 6 — Visual Reflection: Interactive color/shape choices capturing non-verbal emotional state.",
        "Step 13 — Reflection Journal: Distraction-free writing canvas with guided prompt cards ('What made you smile today?').",
        "RLS Protection: Journal entries are encrypted in Postgres and unreadable by platform administrators."
    ], "section3_4_journal.png", "Reflection Journal — serene, prompt-guided writing canvas", "right_text")

    # SLIDE 8: PERSONAL ANALYTICS
    add_dual_image_slide("PERSONAL ANALYTICS", "Insights & PDF Reports",
        ("section3_11_analytics.png", "Insights Page", "Behavioural pattern dashboard powered by Chart.js. Correlates sleep, mood, energy, and productivity across weeks."),
        ("section3_6_feedback.png", "Weekly / Monthly Reports + PDF Export", "Summarises journey activity into weekly/monthly snapshots. One-click jsPDF export compiles charts for personal records or therapy.")
    )

    # SLIDE 9: AI COMPANION
    add_single_image_slide("AI FEATURE", "AI Companion — Lighthouse Chat", [
        "Context-Aware AI: Powered by Groq (Llama-3.3-70B) via an Express.js backend proxy.",
        "Personal Data Grounding: Reads the user's real check-in averages and mood trends to give grounded, empathetic guidance.",
        "Privacy Guardrails: Never receives raw journal text; receives only aggregated numerical metrics."
    ], "section3_10_companion.png", "Lighthouse AI Companion — empathetic, grounded wellness guidance", "left_text")

    # SLIDE 10: ADMIN PORTAL - MONITORING & CARE
    add_dual_image_slide("ADMIN PORTAL", "Platform Monitoring & Care Panel",
        ("section3_12_admin_analytics.png", "Admin Dashboard & Analytics", "Bird's-eye view of platform health: total users, daily activity, growth charts, flagged counts, and PDF report packs."),
        ("section3_14_care_panel.png", "Care Panel — Watchlist & Nudges", "Flagged users appear here via heuristics (low mood streaks, inactivity). Admins review signals, add notes, or send soft nudges.")
    )

    # SLIDE 11: ADMIN PORTAL - USER DIRECTORY & SCENARIOS
    add_dual_image_slide("ADMIN PORTAL", "User Directory & Scenario Management",
        ("section3_13_user_directory.png", "User Directory", "Admins browse all registered users, view engagement status, and surface accounts needing care with zero journal access."),
        ("section3_15_scenarios.png", "Scenario Bank Management", "Admins create, edit, enable, or disable scenario items. Ships with 80+ seed scenarios customizable to institutional context.")
    )

    # SLIDE 12: TECH STACK
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "ARCHITECTURE", "Technology Stack")
    stacks = [
        ("Vanilla HTML / CSS / JS", "Frontend", "Static multi-page app — intentionally no framework. Fast, zero-build, deployable anywhere."),
        ("Supabase (Postgres + RLS)", "Database & Auth", "Row Level Security enforces privacy boundaries at the DB layer — users only see their own rows."),
        ("Groq · Llama-3.3-70B", "AI Backend", "Express.js proxy server holds the Groq API key server-side, feeding live data snapshots per request."),
        ("Chart.js", "Data Visualisation", "Renders all mood, sleep, energy, productivity, and behavioural trend charts on dashboard & insights."),
        ("jsPDF (client-side)", "PDF Export", "Generates personal wellness reports and admin PDF packs entirely in the browser."),
        ("Vercel (Static Deploy)", "Hosting", "Deployed via vercel.json. Zero-build static deployment with environment configuration via config.js.")
    ]
    for idx, (n, l, d) in enumerate(stacks):
        col = idx % 3
        row = idx // 3
        lx = Inches(0.8 + col * 3.9)
        ty = Inches(1.6 + row * 2.7)
        
        c = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, ty, Inches(3.7), Inches(2.4))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(255, 255, 255)
        c.line.color.rgb = BORDER_BLUE
        
        tb = s12.shapes.add_textbox(lx + Inches(0.2), ty + Inches(0.2), Inches(3.3), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = l.upper()
        p0.font.name = 'Segoe UI'
        p0.font.size = Pt(9.5)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT_BLUE
        p0.space_after = Pt(2)
        
        p1 = tf.add_paragraph()
        p1.text = n
        p1.font.name = 'Segoe UI'
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_MAIN
        p1.space_after = Pt(4)
        
        p2 = tf.add_paragraph()
        p2.text = d
        p2.font.name = 'Segoe UI'
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    # SLIDE 13: CONCLUSION
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "CONCLUSION", "What Lighthouse Delivers")
    concls = [
        ("Guided Daily Habit", "A 13-step journey building consistent self-reflection — check-ins, scenarios, activities, visual, journal."),
        ("Bulletproof Privacy", "Supabase RLS means journal entries are cryptographically scoped to the user. Admins see signals, never diaries."),
        ("Personal Insights", "Longitudinal charts surface mood-sleep-energy correlations that single-point wellness checks miss."),
        ("Compassionate Admin Care", "Heuristic flagging + soft nudges let admins reach out proactively — as a care tool, not surveillance."),
        ("AI Companion", "Context-aware Groq-powered assistant grounded in real user data — no hallucinated scores or generic advice."),
        ("Future Roadmap", "Component framework migration, mobile app, and predictive anomaly detection for admin care signals.")
    ]
    for idx, (h, d) in enumerate(concls):
        col = idx % 2
        row = idx // 3
        lx = Inches(0.8) if col == 0 else Inches(6.9)
        ty = Inches(1.6 + (idx % 3) * 1.75)
        
        c = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, ty, Inches(5.6), Inches(1.55))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(255, 255, 255)
        c.line.color.rgb = BORDER_BLUE
        
        tb = s13.shapes.add_textbox(lx + Inches(0.3), ty + Inches(0.15), Inches(5.0), Inches(1.25))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = h
        p1.font.name = 'Segoe UI'
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_BLUE
        p1.space_after = Pt(2)
        p2 = tf.add_paragraph()
        p2.text = d
        p2.font.name = 'Segoe UI'
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    # SLIDE 14: THANK YOU
    s14 = prs.slides.add_slide(blank_layout)
    add_bg(s14)
    
    logo14 = s14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.16), Inches(1.2), Inches(1.0), Inches(1.0))
    logo14.fill.solid()
    logo14.fill.fore_color.rgb = ACCENT_BLUE
    logo14.line.fill.background()
    
    tb_ty = s14.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(1.0))
    p = tb_ty.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    tb_q = s14.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(11.333), Inches(0.4))
    p = tb_q.text_frame.paragraphs[0]
    p.text = "Questions & Discussion"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED
    
    ty_card = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.4), Inches(8.333), Inches(2.0))
    ty_card.fill.solid()
    ty_card.fill.fore_color.rgb = BLUE_BG
    ty_card.line.color.rgb = BORDER_BLUE
    
    tb_ty_info = s14.shapes.add_textbox(Inches(2.7), Inches(4.6), Inches(7.933), Inches(1.6))
    tf_ty = tb_ty_info.text_frame
    tf_ty.word_wrap = True
    
    p_proj = tf_ty.paragraphs[0]
    p_proj.text = "PROJECT LIGHTHOUSE"
    p_proj.alignment = PP_ALIGN.CENTER
    p_proj.font.name = 'Segoe UI'
    p_proj.font.size = Pt(11)
    p_proj.font.bold = True
    p_proj.font.color.rgb = ACCENT_BLUE
    p_proj.space_after = Pt(6)
    
    p_team = tf_ty.add_paragraph()
    p_team.text = "Muhammad Zarif Rahman  ·  Md. Al-Saiban  ·  Md. Irfan"
    p_team.alignment = PP_ALIGN.CENTER
    p_team.font.name = 'Segoe UI'
    p_team.font.size = Pt(14)
    p_team.font.bold = True
    p_team.font.color.rgb = TEXT_MAIN
    p_team.space_after = Pt(4)
    
    p_course = tf_ty.add_paragraph()
    p_course.text = "Software Development with JAVA (Sessional) (CSE-202)  ·  CUET"
    p_course.alignment = PP_ALIGN.CENTER
    p_course.font.name = 'Segoe UI'
    p_course.font.size = Pt(12)
    p_course.font.color.rgb = TEXT_MUTED

    prs.save("Lighthouse_Presentation.pptx")
    print("Successfully generated Lighthouse_Presentation.pptx with all 14 slides!")

if __name__ == '__main__':
    create_deck()
